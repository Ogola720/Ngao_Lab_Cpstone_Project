import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation(output_path):
    prs = Presentation()
    
    # 1. Title Slide
    slide_layout = prs.slide_layouts[0] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Agricultural Micro-Loan Default Prediction"
    subtitle.text = "Ngao Labs — Capstone Project 2026\nAn end-to-end Machine Learning & Explainable AI solution"

    # 2. Agenda
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Presentation Agenda"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Problem & Context"
    p = tf.add_paragraph()
    p.text = "Data & Methodology"
    p = tf.add_paragraph()
    p.text = "Model Architectures"
    p = tf.add_paragraph()
    p.text = "Performance Results"
    p = tf.add_paragraph()
    p.text = "Web Application"
    p = tf.add_paragraph()
    p.text = "Responsible AI & Findings"

    # 3. Problem Statement
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Context: The Problem"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Kenyan microfinance institutions face significant risk from defaults."
    tf.add_paragraph().text = "Manual credit assessments are slow, inconsistent, and prone to bias."
    tf.add_paragraph().text = "~21.8% of loans default, causing huge financial losses."
    tf.add_paragraph().text = "Smallholder farmers lack formal credit histories."
    tf.add_paragraph().text = "Value Propositions:"
    p = tf.add_paragraph()
    p.text = "1. Faster loan processing (real-time risk scoring)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Reduced default rates (data-driven decisions)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Transparent lending (Explainable AI)"
    p.level = 1

    # 4. Dataset Overview
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only
    title = slide.shapes.title
    title.text = "Data: Dataset Overview"
    
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(5.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Three datasets linked by customer ID:"
    tf.add_paragraph().text = "1. traindemographics (4,346 records)\n   Age, employment, education, bank"
    tf.add_paragraph().text = "2. trainprevloans (18,183 records)\n   Historical loans, amounts, dates"
    tf.add_paragraph().text = "3. trainperf (4,368 records)\n   Current loan performance + target"
    tf.add_paragraph().text = "\nTarget: good_bad_flag"
    tf.add_paragraph().text = "Class Distribution: 78.2% Repaid vs 21.8% Defaulted"
    
    img_path = "Data distribution and Correlation.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(5.0), Inches(2.0), width=Inches(4.5))

    # 5. Feature Engineering
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Methodology: Feature Engineering"
    tf = slide.placeholders[1].text_frame
    tf.text = "Behavioral Features:"
    p = tf.add_paragraph()
    p.text = "Repayment delay days, late payment indicators, interest accrued, loan intensity"
    p.level = 1
    tf.add_paragraph().text = "Per-Customer Aggregations:"
    p = tf.add_paragraph()
    p.text = "Total previous loans, avg repayment delay, on-time rate"
    p.level = 1
    tf.add_paragraph().text = "Demographic Derivatives:"
    p = tf.add_paragraph()
    p.text = "Age at loan application, target encoded bank names"
    p.level = 1
    tf.add_paragraph().text = "Preprocessing Pipeline:"
    p = tf.add_paragraph()
    p.text = "StandardScaler, TargetEncoder, OneHotEncoder"
    p.level = 1

    # 6. Handling Class Imbalance
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Methodology: Handling Class Imbalance"
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.0))
    tf = txBox.text_frame
    tf.text = "Strategies employed to handle 78% vs 22% imbalance:"
    tf.add_paragraph().text = "1. SMOTE Oversampling\n   Generates synthetic default examples."
    tf.add_paragraph().text = "2. Cost-Sensitive Weighting\n   Penalizes misclassifying defaults more heavily."
    tf.add_paragraph().text = "3. Threshold Optimization\n   Shifted threshold from 0.5 to ~0.25 to maximize default recall."
    
    img_path = "Confusion Martices with SMOTE.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(5.0), Inches(2.0), width=Inches(4.5))

    # 7. Model Architectures
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Models: Architectures"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. XGBoost Tuned (Primary / Production Model)"
    p = tf.add_paragraph()
    p.text = "Gradient-boosted tree ensemble (Val AUC-ROC: 0.724)"
    p.level = 1
    tf.add_paragraph().text = "2. PyTorch DeepTabNet"
    p = tf.add_paragraph()
    p.text = "Deep tabular comparison model (Val AUC-ROC: 0.716)"
    p.level = 1
    tf.add_paragraph().text = "3. PyTorch ShallowMLP"
    p = tf.add_paragraph()
    p.text = "Neural network benchmark (Val AUC-ROC: 0.694)"
    p.level = 1
    
    # 8. End-to-End Pipeline
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Architecture: End-to-End Pipeline"
    
    img_path = "architecture_diagram_1786022620401.jpg"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.0), Inches(1.5), width=Inches(8.0))

    # 9. Model Performance (AUC-ROC)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Results: Model Performance (AUC-ROC)"
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.0), Inches(5.0))
    tf = txBox.text_frame
    tf.text = "XGBoost consistently outperforms both neural network architectures on this tabular dataset."
    tf.add_paragraph().text = "\nBest Model: Tuned XGBoost (0.724)"
    
    img_path = "AUC & ROC Curves.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(4.5), Inches(1.5), width=Inches(5.0))

    # 10. Threshold Tuning
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Results: Threshold Tuning & F1-Score"
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.0), Inches(5.0))
    tf = txBox.text_frame
    tf.text = "Key Insight:"
    tf.add_paragraph().text = "Shifting the XGBoost threshold from 0.5 to ~0.25 was more effective at catching defaults than SMOTE alone."
    
    img_path = "Comparative F1 scores with tuned XGboost Threshold.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(4.5), Inches(1.5), width=Inches(5.0))
    elif os.path.exists("Comparative F1 score with SMOTE.png"):
        slide.shapes.add_picture("Comparative F1 score with SMOTE.png", Inches(4.5), Inches(1.5), width=Inches(5.0))

    # 11. Confusion Matrices
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Results: Confusion Matrices (Tuned)"
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.0), Inches(5.0))
    tf = txBox.text_frame
    tf.text = "After threshold tuning:"
    tf.add_paragraph().text = "- Dramatically improved default recall"
    tf.add_paragraph().text = "- Precision vs recall balanced per business need"
    
    img_path = "Confusion Matrixes with tuned XGboost threshold.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(4.5), Inches(1.5), width=Inches(5.0))
    elif os.path.exists("Confusion Matrixes.png"):
        slide.shapes.add_picture("Confusion Matrixes.png", Inches(4.5), Inches(1.5), width=Inches(5.0))

    # 12. Web Application
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Deployment: AgriLoan Predictor Web App"
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.0), Inches(5.0))
    tf = txBox.text_frame
    tf.text = "Interactive Streamlit dashboard for real-time loan risk assessment."
    tf.add_paragraph().text = "- Loan Input Form (Demographics, Banking, History)"
    tf.add_paragraph().text = "- Real-Time Prediction (High/Low Risk)"
    tf.add_paragraph().text = "- SHAP Explanations (Top 10 features)"
    
    # use first available webapp screenshot
    for i in range(1, 7):
        img_path = f"Webapp screenshot_{i}.png"
        if not os.path.exists(img_path):
            img_path = f"Webapp_screenshot_{i}.png"
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(4.5), Inches(1.5), width=Inches(5.0))
            break

    # 13. Responsible AI
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Ethics: Responsible AI Framework"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. Explainability (SHAP)"
    p = tf.add_paragraph()
    p.text = "Every prediction includes feature impact analysis, not a black box."
    p.level = 1
    tf.add_paragraph().text = "2. Fairness Audit"
    p = tf.add_paragraph()
    p.text = "Consistent performance across employment sub-groups."
    p.level = 1
    tf.add_paragraph().text = "3. Transparency Disclaimer"
    p = tf.add_paragraph()
    p.text = "AI is a decision-support tool, not a replacement for human judgment."
    p.level = 1

    # 14. Key Findings
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Insights: Key Findings"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. Top Risk Drivers"
    p = tf.add_paragraph()
    p.text = "Repayment delays and high loan intensity are the strongest predictors."
    p.level = 1
    tf.add_paragraph().text = "2. Threshold Tuning > SMOTE"
    p = tf.add_paragraph()
    p.text = "Adjusting classification threshold to ~0.25 was most effective."
    p.level = 1
    tf.add_paragraph().text = "3. XGBoost > Neural Networks"
    p = tf.add_paragraph()
    p.text = "Gradient boosting remains king for structured data."
    p.level = 1
    tf.add_paragraph().text = "4. Feature Engineering Matters"
    p = tf.add_paragraph()
    p.text = "Behavioral features had far higher predictive power than demographics alone."
    p.level = 1

    # 15. Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Thank You"
    subtitle.text = "Built with ❤️ for Kenyan Agricultural Finance\nEmpowering farmers, one loan at a time 🌾"
    
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == '__main__':
    create_presentation('Agricultural_Micro_Loan_Default_Prediction.pptx')
